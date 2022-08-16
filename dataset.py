import numpy as np
import torch
from pycocotools.coco import COCO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE
from tqdm import tqdm
from torchvision.datasets.mnist import MNIST
from torch.utils.data.dataset import Dataset
from PIL import Image, ImageDraw, ImageOps
import seaborn as sns
import json
import os

import random
import pdb

from utils import trim_tokens, gen_colors

def get_dataset(name, split, max_length=None, precision=8):
    if name == "rico":
        return RicoLayout(split, max_length, precision)
    elif name == "publaynet":
        return PubLayout(split, max_length, precision)
    elif name == "infoppt":
        return PPTLayout(split, max_length, precision)
    
    raise NotImplementedError(name)

class Padding(object):
    def __init__(self, max_length, vocab_size):
        self.max_length = max_length
        self.bos_token = vocab_size - 3
        self.eos_token = vocab_size - 2
        self.pad_token = vocab_size - 1

    def __call__(self, layout):
        # grab a chunk of (max_length + 1) from the layout

        chunk = torch.zeros(self.max_length+1, dtype=torch.long) + self.pad_token
        # Assume len(item) will always be <= self.max_length:
        chunk[0] = self.bos_token
        chunk[1:len(layout)+1] = layout
        chunk[len(layout)+1] = self.eos_token

        x = chunk[:-1]
        y = chunk[1:]
        return {'x': x, 'y': y}

class Mask_and_padding_BLT(object):
    def __init__(self, max_length, vocab_size):
        self.max_length = max_length
        self.mask_id = vocab_size - 4
        self.bos_token = vocab_size - 3
        self.eos_token = vocab_size - 2
        self.pad_token = vocab_size - 1
        # self.masked_lm_prob = masked_lm_prob

    def __call__(self, tokens, element_length):
        seq_length = len(tokens)

        '''
        Mask:
        - C,S,P grouped mask
        '''
        random_group = random.randint(0,2)
        if random_group == 0: #C
            num_predictions = random.randint(1, element_length)
            positions_sample = list(random.sample(range(element_length), num_predictions)) ##### random ways
            positions_sample.sort()
            masked_lm_positions = np.array([5 * i for i in positions_sample])

            # print("group: Category, masked token:",self.mask_id,", masked_num:",num_predictions)

        else:
            num_predictions = random.randint(1, element_length * 2)
            positions_sample = list(random.sample(range(element_length * 2), num_predictions))  ##### random ways
            positions_sample.sort()
            if random_group == 1:
                masked_lm_positions = np.array([5 * (i // 2) + i % 2 + 1 for i in positions_sample])
                # print("group: Position, masked token:", self.mask_id, ", masked_num:", num_predictions)
            else:
                masked_lm_positions = np.array([5 * (i // 2) + i % 2 + 3 for i in positions_sample])
                # print("group: Size, masked token:", self.mask_id, ", masked_num:", num_predictions)


        masked_tokens = np.array(tokens).copy()
        masked_tokens[masked_lm_positions] = self.mask_id
        masked_tokens = torch.tensor(masked_tokens,dtype=torch.long)
        tokens = torch.tensor(tokens,dtype=torch.long)

        masked_seq = np.full(seq_length,False)
        masked_seq[masked_lm_positions] = True
        masked_seq = torch.tensor(masked_seq,dtype=torch.bool)


        '''
        Padding:
        - pad into target length
        '''
        chunk_mask = torch.zeros(self.max_length+1, dtype=torch.long) + self.pad_token
        masked_seq_tmp = torch.full((chunk_mask.size()),False)
        
        # Assume len(item) will always be <= self.max_length:
        chunk_mask[0] = self.bos_token
        chunk_mask[1:seq_length+1] = masked_tokens
        x = chunk_mask[:-1]
        masked_seq_tmp[1:seq_length+1] = masked_seq
        masked_seq_x = masked_seq_tmp[:-1]
        masked_seq_y = masked_seq_tmp[1:]

        #### original output: tokens + eos +  padding
        chunk = torch.zeros(self.max_length+1, dtype=torch.long) + self.pad_token
        chunk[1:seq_length+1] = tokens
        chunk[seq_length + 2] = self.eos_token
        y = chunk[1:]

        chunk_anti_mask_output = torch.zeros(self.max_length+1, dtype=torch.long) # + self.pad_token
        anti_mask_tokens = masked_seq.float() * tokens
        chunk_anti_mask_output[1:seq_length + 1] = anti_mask_tokens
        # chunk_anti_mask_output[seq_length + 2] = self.eos_token
        lm_output = chunk_anti_mask_output[1:]
        # return {'x': x, 'y': y}
        # print(tokens)
        # print("x",x)
        # print("y",y)
        # print(lm_output)
        # pdb.set_trace()
        return x,y,lm_output,masked_seq_x,masked_seq_y,torch.tensor(seq_length,dtype=torch.int64)


class MNISTLayout(MNIST):

    def __init__(self, root, train=True, download=True, threshold=32, max_length=None):
        super().__init__(root, train=train, download=download)
        self.vocab_size = 784 + 3  # bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1

        self.threshold = threshold
        self.data = [self.img_to_set(img) for img in self.data]
        self.max_length = max_length
        if self.max_length is None:
            self.max_length = max([len(x) for x in self.data]) + 2  # bos, eos tokens
        self.transform = Padding(self.max_length, self.vocab_size)

    def __len__(self):
        return len(self.data)

    def img_to_set(self, img):
        fg_mask = img >= self.threshold
        fg_idx = fg_mask.nonzero(as_tuple=False)
        fg_idx = fg_idx[:, 0] * 28 + fg_idx[:, 1]
        return fg_idx

    def render(self, layout):
        layout = trim_tokens(layout, self.bos_token, self.eos_token, self.pad_token)
        x_coords = layout % 28
        y_coords = layout // 28
        # valid_idx = torch.where((y_coords < 28) & (y_coords >= 0))[0]
        img = np.zeros((28, 28, 3)).astype(np.uint8)
        img[y_coords, x_coords] = 255
        return Image.fromarray(img, 'RGB')

    def __getitem__(self, idx):
        # grab a chunk of (block_size + 1) tokens from the data
        layout = self.transform(self.data[idx])
        return layout['x'], layout['y']

class JSONLayout(Dataset):
    def __init__(self, json_path, max_length=None, precision=8):
        with open(json_path, "r") as f:
            data = json.loads(f.read())

        images, annotations, categories = data['images'], data['annotations'], data['categories']
        self.size = pow(2, precision)
        #fatgoose 255 -> 31

        self.categories = {c["id"]: c for c in categories}
        self.colors = gen_colors(len(self.categories))

        self.json_category_id_to_contiguous_id = {
            v: i + self.size for i, v in enumerate([c["id"] for c in self.categories.values()])
        }
        # {1: 256, 2: 257, 3: 258, 4: 259, 5: 260}

        self.contiguous_category_id_to_json_id = {
            v: k for k, v in self.json_category_id_to_contiguous_id.items()
        }

        self.vocab_size = self.size + len(self.categories) + 3  # bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1

        image_to_annotations = {}
        for annotation in annotations:
            image_id = annotation["image_id"]

            if not (image_id in image_to_annotations):
                image_to_annotations[image_id] = []

            image_to_annotations[image_id].append(annotation)

        self.data = []
        for image in images:
            image_id = image["id"]
            height, width = float(image["height"]), float(image["width"])

            if image_id not in image_to_annotations:
                continue

            ann_box = []
            ann_cat = []

            for ann in image_to_annotations[image_id]:
                x, y, w, h = ann["bbox"]
                ann_box.append([x, y, w, h])
                ann_cat.append(self.json_category_id_to_contiguous_id[ann["category_id"]])

            # Sort boxes
            ann_box = np.array(ann_box)
            ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
            # Sort by ann_box[:, 1], then by ann_box[:, 0]
            ann_box = ann_box[ind]
       
            if len(ann_box) > 128:
                print("elements num: ", len(ann_box))
                continue

            ann_cat = np.array(ann_cat)
            ann_cat = ann_cat[ind]

            # Discretize boxes
            ann_box = self.quantize_box(ann_box, width, height)

            # Append the categories
            layout = np.concatenate([ann_cat.reshape(-1, 1), ann_box], axis=1)
            # Flatten and add to the dataset
            self.data.append(layout.reshape(-1))

        self.max_length = max_length
        if self.max_length is None:
            self.max_length = max([len(x) for x in self.data]) + 2  # bos, eos tokens
        self.transform = Padding(self.max_length, self.vocab_size)

    def quantize_box(self, boxes, width, height):

        # range of xy is [0, large_side-1]
        # range of wh is [1, large_side]
        # bring xywh to [0, 1]
        boxes[:, [2, 3]] = boxes[:, [2, 3]] - 1
        boxes[:, [0, 2]] = boxes[:, [0, 2]] / (width - 1)
        boxes[:, [1, 3]] = boxes[:, [1, 3]] / (height - 1)
        boxes = np.clip(boxes, 0, 1)

        # next take xywh to [0, size-1]
        boxes = (boxes * (self.size - 1)).round()
        # round: 四舍五入

        return boxes.astype(np.int32)

    def __len__(self):
        return len(self.data)

    def render(self, layout):
        img = Image.new('RGB', (256, 256), color=(255, 255, 255))
        draw = ImageDraw.Draw(img, 'RGBA')
        layout = layout.reshape(-1)
        layout = trim_tokens(layout, self.bos_token, self.eos_token, self.pad_token)
        layout = layout[: len(layout) // 5 * 5].reshape(-1, 5)
        box = layout[:, 1:].astype(np.float32)
        box[:, [0, 1]] = box[:, [0, 1]] / (self.size - 1) * 255
        box[:, [2, 3]] = box[:, [2, 3]] / self.size * 256
        box[:, [2, 3]] = box[:, [0, 1]] + box[:, [2, 3]]

        for i in range(len(layout)):
            x1, y1, x2, y2 = box[i]
            cat = layout[i][0]
            col = self.colors[cat-self.size] if 0 <= cat-self.size < len(self.colors) else [0, 0, 0]
            draw.rectangle([x1, y1, x2, y2],
                           outline=tuple(col) + (200,),
                           fill=tuple(col) + (64,),
                           width=2)

        # Add border around image
        img = ImageOps.expand(img, border=2)
        return img

    def __getitem__(self, idx):
        # grab a chunk of (block_size + 1) tokens from the data
        layout = torch.tensor(self.data[idx], dtype=torch.long)
        layout = self.transform(layout)
        return layout['x'], layout['y']

class BaseDataset(Dataset):
    component_class = []
    _category_id_to_category_name = None
    _json_category_id_to_contiguous_id = None
    _contiguous_id_to_json_id = None
    _colors = None

    def __init__(self, name, split, is_rela=False):
        super().__init__()
        assert split in ['train', 'val', 'test']
        self.is_rela = is_rela
        dir_path = os.path.dirname(os.path.realpath(__file__))
        idx = self.processed_file_names.index('{}.pt'.format(split))
        os.makedirs(os.path.join(dir_path, "preprocess_data", name), exist_ok=True)
        self.data_path = os.path.join(dir_path, "preprocess_data", name, self.processed_file_names[idx])
        self.W = 256
        self.H = 256


    @property
    def json_category_id_to_contiguous_id(self):
        if self._json_category_id_to_contiguous_id is None:
            self._json_category_id_to_contiguous_id = {
            i: i + self.size + 1 for i in range(self.categories_num)
        }
        return self._json_category_id_to_contiguous_id

    @property
    def contiguous_category_id_to_json_id(self):
        if self._contiguous_id_to_json_id is None:
            self._contiguous_category_id_to_json_id = {
            v: k for k, v in self.json_category_id_to_contiguous_id.items()
        }
        return self._contiguous_id_to_json_id

    @property
    def colors(self):
        if self._colors is None:
            num_colors = self.categories_num
            palette = sns.color_palette(None, num_colors)
            if num_colors > 10:
                palette[10:] = sns.color_palette("husl", num_colors-10)
            self._colors = [[int(x[0]*255), int(x[1]*255), int(x[2]*255)] for x in palette]
        return self._colors

    @property
    def processed_file_names(self):
        return ['train.pt', 'val.pt', 'test.pt']

    def quantize_box(self, boxes, width, height):
        # range of xy is [0, large_side-1]
        # range of wh is [1, large_side]
        # bring xywh to [0, 1]
        
        boxes[:, [2, 3]] = boxes[:, [2, 3]] - 1
        boxes[:, [0, 2]] = boxes[:, [0, 2]] / (width - 1)
        boxes[:, [1, 3]] = boxes[:, [1, 3]] / (height - 1)
        boxes = np.clip(boxes, 0, 1)

        # next take xywh to [0, size-1]
        boxes = (boxes * (self.size - 1)).round()
        # round: 四舍五入

        return boxes.astype(np.int32)

    def __len__(self):
        return len(self.data)

    def render_normalized_layout(self, layout):
        layout = layout.reshape(-1)
        layout = trim_tokens(layout, self.bos_token, self.eos_token, self.pad_token)
        layout = layout[: len(layout) // 5 * 5].reshape(-1, 5)
        box = layout[:, 1:].astype(np.float32)
        label = layout[:, 0].astype(np.int32)
        label = label - self.size
        box = box / (self.size - 1)
        return (box, label)

    def render(self, layout):
        img = Image.new('RGB', (self.W, self.H), color=(255, 255, 255))
        draw = ImageDraw.Draw(img, 'RGBA')
        layout = layout.reshape(-1)
        layout = trim_tokens(layout, self.bos_token, self.eos_token, self.pad_token)
        layout = layout[: len(layout) // 5 * 5].reshape(-1, 5)
        box = layout[:, 1:].astype(np.float32)
        # box[:, [0, 1]] = box[:, [0, 1]] / (self.size - 1) * 255
        # box[:, [2, 3]] = box[:, [2, 3]] / self.size * 256
        # box[:, [2, 3]] = box[:, [0, 1]] + box[:, [2, 3]]
        box = box / (self.size - 1)

        box[:, [0, 2]] = box[:, [0, 2]] * self.W
        box[:, [1, 3]] = box[:, [1, 3]] * self.H
        # xywh to ltrb
        x1s = box[:, 0] - box[:, 2] / 2
        y1s = box[:, 1] - box[:, 3] / 2
        x2s = box[:, 0] + box[:, 2] / 2
        y2s = box[:, 1] + box[:, 3] / 2

        for i in range(len(layout)):
            # x1, y1, x2, y2 = box
            x1, y1, x2, y2 = x1s[i], y1s[i], x2s[i], y2s[i]
            cat = layout[i][0]
            col = self.colors[cat-self.size] if 0 <= cat-self.size < len(self.colors) else [0, 0, 0]
            draw.rectangle([x1, y1, x2, y2],
                           outline=tuple(col) + (200,),
                           fill=tuple(col) + (64,),
                           )

        # Add border around image
        img = ImageOps.expand(img, border=2)
        return img

    def __getitem__(self, idx):
        # grab a chunk of (block_size + 1) tokens from the data
        layout = torch.tensor(self.data[idx], dtype=torch.long)
        layout = self.transform(layout)
        return layout['x'], layout['y']  

    def load_pt(self, load_path):
        results = torch.load(load_path)
        self.categories_num = results["categories_num"]
        self.max_elements_num = results["max_elements_num"]
        self.data = results["data"]

        if "iou_data" in results:
            print("load iou data")
            self.iou_data = results["iou_data"]

        self.copy_mode = self.size + 1 + self.categories_num
        self.margin_mode = self.copy_mode + 1
        self.generate_mode = self.margin_mode + 1
        self.no_obj_token = self.generate_mode + 1

class RicoLayout(BaseDataset):
    component_class = {'Toolbar':0, 'Image':1, 'Text':2, 'Icon':3, 'Text Button':4, 'Input':5,
        'List Item': 6, 'Advertisement': 7, 'Pager Indicator':8, 'Web View':9, 'Background Image':10,
        'Drawer':11, 'Modal':12}
    def __init__(self, split, max_length=None, precision=8):    
        super().__init__('rico', split)
        # component_class = {'Text':0, 'Icon':1, 'Image':2, 'Text Button':3, 'Toolbar':4, 'List Item':5, 'Web View':6, 
        # 'Advertisement':7, 'Input':8, 'Drawer':9, 'Background Image':10, 'Card':11, 'Multi-Tab':12, 'Modal':13, 
        # 'Pager Indicator':14, 'Radio Button':15, 'On/Off Switch':16, 'Slider':17, 'Checkbox':18, 'Map View':19,
        # 'Button Bar':20, 'Video':21, 'Bottom Navigation':22, 'Date Picker':23, 'Number Stepper':24}
        self.categories_num = len(self.component_class.keys())
        self.size = pow(2, precision)
        self.vocab_size = self.size + self.categories_num + 4  # bos, eos, pad, mask tokens
        self.mask_id = self.vocab_size - 4
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1

        if os.path.exists(self.data_path):
            print("load dataset.")
            self.load_pt(self.data_path)
        else:
            self.data = []
            self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
            bbox_idx = 0
            self.max_elements_num = 9
            data_dir = f"/home/v-jierulin/layoutTransformer/LayoutGeneration/DeepLayout/layout_transformer_BLT/dataset/rico/semantic_annotations"
            dirs = os.listdir(data_dir)         
            for file in dirs:
                if file.split(".")[-1] == "json":
                    file_path = os.path.join(data_dir, file)
                    with open(file_path, "r") as f:
                        json_file = json.load(f)

                    canvas = json_file["bounds"]
                    W, H = float(canvas[2]-canvas[0]), float(canvas[3]-canvas[1])
                    if canvas[0]!= 0 or canvas[1]!= 0 or W <= 1000:
                        continue
                    elements = self.get_all_element(json_file, [])
                    elements = list(filter(lambda e: e["componentLabel"] in self.component_class, elements))
                    
                    
                    if len(elements) == 0 or len(elements)>self.max_elements_num:
                        continue

                    ann_box = []
                    ann_cat = []

                    for ele in elements:
                        [x_l, y_t, x_r, y_b] = ele["bounds"]
                        xc = (x_l + x_r) / 2.
                        yc = (y_t + y_b) / 2.
                        w = x_r - x_l
                        h = y_b - y_t

                        if w<0 or h<0:
                            continue
                        ann_box.append([xc, yc, w, h])
                        ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[ele["componentLabel"]]])

                    # Sort boxes

                    ann_box = np.array(ann_box)
                    # Discretize boxes
                    ann_box = self.quantize_box(ann_box, W, H)

                    ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                    # Sort by ann_box[:, 1], then by ann_box[:, 0]
                    ann_box = ann_box[ind]
                    
                    ann_cat = np.array(ann_cat)
                    ann_cat = ann_cat[ind]

                    self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                    self.iou_data["file_idx"].append(int(file.split(".")[0]))
                    self.iou_data["file2bboxidx"][int(file.split(".")[0])] = bbox_idx

                    bbox_idx += 1

                    # Append the categories
                    layout = np.concatenate([ann_cat.reshape(-1, 1), ann_box], axis=1)

                    # Flatten and add to the dataset
                    self.data.append(layout.reshape(-1))

            self.save_pt(self.data_path)  

        # pdb.set_trace() # try to look at dataset raw


        self.max_length = max_length
        if self.max_length is None:
            self.max_length = max([len(x) for x in self.data]) + 2  # bos, eos tokens

        # self.transform = Padding(self.max_length, self.vocab_size)
        self.transform = Mask_and_padding_BLT(self.max_length, self.vocab_size)

        self.masked_lm_prob = 0.5

    
    def build_training_sample(self, tokens):
        """Biuld training sample.
        Arguments:
            sample: A list of sentences in which each sentence is a list token ids.
            target_seq_length: Desired sequence length.
            max_seq_length: Maximum length of the sequence. All values are padded to
                this length.
            vocab_id_list: List of vocabulary ids. Used to pick a random id.
            vocab_id_to_token_dict: A dictionary from vocab ids to text tokens.
            cls_id: Start of example id.
            sep_id: Separator id.
            mask_id: Mask token id.
            pad_id: Padding token id.
            masked_lm_prob: Probability to mask tokens.
            np_rng: Random number genenrator. Note that this rng state should be
                numpy and not python since python randint is inclusive for
                the opper bound whereas the numpy one is exclusive.
        """
        # assert target_seq_length <= max_seq_length

        # tokens_np, tokentypes_np, labels_np, padding_mask_np, loss_mask_np = self.transform(tokens, vocab_id_list, vocab_id_to_token_dict,mask_id, np_rng)
        # pdb.set_trace()
        assert len(tokens) % 5 == 0
        x,y,lm_output, masked_seq_x,masked_seq_y,element_num = self.transform(tokens,int(len(tokens)/5))
        # # Masking. #fatgoose
        # (tokens, masked_positions, masked_labels, _, _) = create_masked_lm_predictions(tokens, vocab_id_list, vocab_id_to_token_dict, masked_lm_prob,cls_id, sep_id, mask_id, max_predictions_per_seq, np_rng)

        # # Padding.
        # # why convert to numpy
        # tokens_np, tokentypes_np, labels_np, padding_mask_np, loss_mask_np = pad_and_convert_to_numpy(tokens, tokentypes, masked_positions,masked_labels, pad_id, max_seq_length)

        train_sample = {
            'text': x,
            'target_output': y,
            'lm_output':lm_output,
            # 'is_random': int(is_next_random),
            'input_mask': masked_seq_x,
            'loss_mask': masked_seq_y,#fatgoose ?
            'element_num':element_num,
            # 'padding_mask': padding_mask_np,
            # 'truncated': int(truncated)
        }

        # pdb.set_trace()

        return {key:torch.tensor(value, dtype=torch.long) for key,value in train_sample.items()}

    def __getitem__(self,idx):
        # grab a chunk of (block_size + 1) tokens from the data
        layout = np.array(self.data[idx])#list(self.data[idx])
        # layout = self.transform(layout)
        # return layout['x'], layout['y']
        # np_rng = np.random.RandomState(seed=((self.seed + idx) % 2**32)) #?
        return self.build_training_sample(layout)

    def get_all_element(self, p_dic, elements):
        if "children" in p_dic:
            for i in range(len(p_dic["children"])):
                cur_child = p_dic["children"][i]
                elements.append(cur_child)
                elements = self.get_all_element(cur_child, elements)
        return elements

    def save_pt(self, save_path):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data))
        s = [int(N * .85), int(N * .90)]
        results["data"] = self.data[:s[0]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[0]))
        results["data"] = self.data[s[0]:s[1]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[1]))
        results["data"] = self.data[s[1]:]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[2]))

class PubLayout(BaseDataset):
    component_class = {'text': 0, 'title':1, 'list':2, 'table':3, 'figure':4}
    def __init__(self, split, max_length=None, precision=8):
        super().__init__('publaynet', split)
        self.categories_num = len(self.component_class.keys())

        self.size = pow(2, precision)
        self.vocab_size = self.size + self.categories_num + 3  # bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1

        self.data = []
        self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
        bbox_idx = 0
        
        if os.path.exists(self.data_path):
            print("load dataset.")
            self.load_pt(self.data_path)
        else:
            ori_split = "train" if split in ["train", "val"] else "val"
            data_path = f"/home/v-jierulin/layoutTransformer/LayoutGeneration/DeepLayout/layout_transformer_BLT/dataset/publaynet/publaynet/{ori_split}.json"
            coco = COCO(data_path)
            self.max_elements_num = 9
            for img_id in sorted(coco.getImgIds()):
                ann_img = coco.loadImgs(img_id)
                W = float(ann_img[0]['width'])
                H = float(ann_img[0]['height'])
                name = ann_img[0]['file_name']
                if H < W:
                    continue

                def is_valid(element):
                    x1, y1, width, height = element['bbox']
                    x2, y2 = x1 + width, y1 + height
                    if x1 < 0 or y1 < 0 or W < x2 or H < y2:
                        return False
                    if x2 <= x1 or y2 <= y1:
                        return False
                    if width <=0 or height <=0:
                        return False
                    return True

                elements = coco.loadAnns(coco.getAnnIds(imgIds=[img_id]))
                elements = list(filter(is_valid, elements))
                
                N = len(elements)
                if N == 0 or self.max_elements_num < N:
                    continue

                ann_box = []
                ann_cat = []

                for element in elements:
                    # bbox
                    x1, y1, width, height = element['bbox']
                    xc = x1 + width / 2.
                    yc = y1 + height / 2.
                    b = [xc , yc , width, height]
                    ann_box.append(b)

                    # label
                    l = coco.cats[element['category_id']]['name']
                    ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[l]])

                # Sort boxes
                ann_box = np.array(ann_box)
                # Sort boxes
                # Discretize boxes
                ann_box = self.quantize_box(ann_box, W, H)

                ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                # Sort by ann_box[:, 1], then by ann_box[:, 0]
                ann_box = ann_box[ind]
                
                ann_cat = np.array(ann_cat)
                ann_cat = ann_cat[ind]

                self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                self.iou_data["file_idx"].append(img_id)
                self.iou_data["file2bboxidx"][img_id] = bbox_idx

                bbox_idx += 1

                # Append the categories
                layout = np.concatenate([ann_cat.reshape(-1, 1), ann_box], axis=1)

                # Flatten and add to the dataset
                self.data.append(layout.reshape(-1))

            self.save_pt(self.data_path, split)

        self.max_length = max_length
        if self.max_length is None:
            self.max_length = max([len(x) for x in self.data]) + 2  # bos, eos tokens
        self.transform = Padding(self.max_length, self.vocab_size)

    def save_pt(self, save_path, split):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data)*0.95)
        if split == "train":
            results["data"] = self.data[:N]
        elif split == "val":
            results["data"] = self.data[N:]
        else:
            results["data"] = self.data
        torch.save(results, save_path)

class PPTLayout(BaseDataset):
    component_class = {"TEXT_BOX": 0, "PICTURE": 1, "CHART":2, "TABLE":3, "TITLE":4, "SUBTITLE":5}   
        
    def __init__(self, split, max_length=None, precision=8):
        super().__init__('infoppt', split)
        SHAPE_TYPE_DICT = {}
        PLACEHOLDER_DICT = {}
        AUTO_SHAPE_DICT = {}
        for idx, name in enumerate(dir(MSO_SHAPE_TYPE), start=1):
            if idx > 26:
                break
            tmp_num = getattr(MSO_SHAPE_TYPE, name)
            SHAPE_TYPE_DICT[tmp_num] = name
        for idx, name in enumerate(dir(PP_PLACEHOLDER), start=1):
            if idx > 22:
                break
            tmp_num = getattr(PP_PLACEHOLDER, name)
            PLACEHOLDER_DICT[tmp_num] = name
        for idx, name in enumerate(dir(MSO_SHAPE), start=1):
            if idx > 184:
                break
            tmp_num = getattr(MSO_SHAPE, name)
            AUTO_SHAPE_DICT[tmp_num] = name

        self.categories_num = len(self.component_class.keys())
        
        # shape.picture + placeholder.picture
        rectangle_class = ['RECTANGLE', 'ROUNDED_RECTANGLE']

        self.categories_num = len(self.component_class.keys())
        
        self.size = pow(2, precision)

        self.vocab_size = self.size + self.categories_num + 3  # bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1

        self.max_element_num = 20
        if os.path.exists(self.data_path):
            print("load InfoPPT dataset.")
            self.load_pt(self.data_path)
        else:
            data_dir = f"../../LayoutAction/datasets/infoppt"
            dirs = os.listdir(data_dir)
            self.data = []
            self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
            bbox_idx = 0
            for file in tqdm(list(dirs), total=len(dirs)): 
                if file.split(".")[-1] == "pptx":
                    file_path = os.path.join(data_dir, file)
                    prs = Presentation(file_path)

                    H, W = prs.slide_height, prs.slide_width   
                    for index, slide in enumerate(prs.slides, start=1):
                        ann_box = []
                        ann_cat = []
                        for shape in slide.shapes:
                            try:
                                # shape.left, shape.top, shape.width, shape.height
                                xc = shape.left + shape.width*0.5
                                yc = shape.top + shape.height*0.5

                                if shape.width <= 0 or shape.height <= 0:
                                    continue
                                
                                if shape.is_placeholder:
                                    cur_class = PLACEHOLDER_DICT[shape.placeholder_format.type]
                                    if cur_class not in ["PICTURE", "TITLE", "SUBTITLE"]:
                                        continue
                                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                    cur_class = AUTO_SHAPE_DICT[shape.auto_shape_type]
                                    if cur_class in rectangle_class and shape.has_text_frame:
                                        if len(shape.text) > 0:
                                            cur_class = "TEXT_BOX"
                                        else:
                                            continue
                                    else:
                                        continue
                                elif SHAPE_TYPE_DICT[shape.shape_type] in self.component_class:
                                    cur_class = SHAPE_TYPE_DICT[shape.shape_type]
                                    if cur_class == "TEXT_BOX" and shape.has_text_frame:
                                        if len(shape.text) == 0:
                                            continue
                                else:
                                    # filter
                                    continue 

                                ann_box.append([xc, yc, shape.width, shape.height])
                                ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[cur_class]])
                            except:
                                pass

                        if len(ann_cat) > self.max_element_num or len(ann_cat) <= 3:
                            continue

                        ann_box = np.array(ann_box)

                        # Sort boxes
                        # Discretize boxes
                        ann_box = self.quantize_box(ann_box, W, H)

                        ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                        # Sort by ann_box[:, 1], then by ann_box[:, 0]
                        ann_box = ann_box[ind]
                        
                        ann_cat = np.array(ann_cat)
                        ann_cat = ann_cat[ind]

                        file_name = file.split(".")[0]
                        self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                        self.iou_data["file_idx"].append(f"{file_name}_{index}")
                        self.iou_data["file2bboxidx"][f"{file_name}_{index}"] = bbox_idx

                        bbox_idx += 1

                        # Append the categories
                        layout = np.concatenate([ann_cat.reshape(-1, 1), ann_box], axis=1)

                        # Flatten and add to the dataset
                        self.data.append(layout.reshape(-1))

            self.save_pt(self.data_path)

        self.max_length = max_length
        if self.max_length is None:
            self.max_length = max([len(x) for x in self.data]) + 2  # bos, eos tokens
        self.transform = Padding(self.max_length, self.vocab_size)


    def save_pt(self, save_path):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data))
        s = [int(N * .85), int(N * .90)]
        results["data"] = self.data[:s[0]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[0]))
        results["data"] = self.data[s[0]:s[1]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[1]))
        results["data"] = self.data[s[1]:]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[2])) 

